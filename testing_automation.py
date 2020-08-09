from os import listdir, getcwd, mkdir, rename, remove
from os.path import isfile, join, exists
from time import strftime
from pptx import Presentation
from pandas import DataFrame
from fpdf import FPDF
from docx import Document
from xlwt import Workbook

timestr = strftime("%Y%m%d-%H%M%S")
path = "C:\\Users\\cawasthi\\Desktop\\WF\\" + timestr
number_of_files = 5

#print(path+"\\ppt")
if  not exists(path):
	mkdir(path)
	mkdir(path+"\\ppt")
	mkdir(path+"\\pdf")
	mkdir(path+"\\txt")
	mkdir(path+"\\docs")
	mkdir(path+"\\excel")	
	print("Directory ", path, " Created ")
else:
	print("Directory ", path, " already exists")
	
#to create PDF
def create_pdf(path):
	pdf = FPDF()
	i = 1
	while i <= number_of_files:	
		pdf = FPDF()
		pdf.add_page()
		pdf.set_font("Arial", size = 15)
		pdf.cell(200, 10, txt = "pdf_doc_number"+timestr+str(i), ln = 1, align = 'C') 
		pdf.output(path+"\\pdf\\"+'pdf_doc_number'+str(i)+timestr+'.pdf')
		i = i + 1

#to create txt
def create_txt(path):
	i = 1
	while i <= number_of_files:
		with open(path+"\\txt\\"+str(i)+timestr+".txt", 'w') as the_file:
			the_file.write("text_number"+str(i)+timestr)
		i = i + 1

#to create xlsx/xls
def create_excel(path):
	i = 1
	while i <= number_of_files:
		book = Workbook(encoding="utf-8")
		sheet1 = book.add_sheet("Sheet 1")
		sheet1.write(0, 0, 'excel_number'+timestr+str(i))
		book.save(path+"\\excel\\"+'excel_number'+str(i)+timestr+'.xls')
		i = i + 1

#to create pptx
def create_presentation(path):
	i = 1
	while i <= number_of_files:
		prs = Presentation() 
		title_slide_layout = prs.slide_layouts[0]
		slide = prs.slides.add_slide(title_slide_layout)
		title = slide.shapes.title
		title.text = 'presentation'+timestr+str(i)
		prs.save(path+"\\ppt\\"+'presentation'+str(i)+timestr+'.pptx')
		i = i + 1
		
#to create docx
def create_document(path):
	i = 1
	while i <= number_of_files:
		document = Document()
		document.add_heading('document'+timestr+str(i), 0)
		document.save(path+"\\docs\\"+'document'+str(i)+timestr+'.docx')
		i = i + 1


#if __name__ == "__main__":
#function_initialize()
create_pdf(path)
create_excel(path)
create_presentation(path)
create_document(path)
create_txt(path)